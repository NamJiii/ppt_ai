#-*- coding:utf-8 -*-
from pptx import Presentation
from pptx.util import Pt
import json

def Start(input):
    prs = Presentation("base.pptx")
    f= open(input, encoding='utf-8')
    return prs, f

#ppt 객체에 새로운 슬라이드를 추가한다.
def NewSlide(prs,slideType,title):
    slide = prs.slides.add_slide(prs.slide_layouts[slideType])#ppt 객체, 슬라이드마스터 번호, 제목
    slide.shapes.title.text = title
    return slide

#제목 슬라이드를 만드는 함수
def TitleSlide(prs,slideType,title,sub_title):#NewSlide, 제목, 부제목
    slide = NewSlide(prs,slideType,title)
    slide.shapes.placeholders[1].text = sub_title
    return slide

#슬라이드 객체의 특정 텍스트박스에 대한 설정
def SetTextBox(slide,cnt,fontAdr):#슬라이드 객체, 텍스트박스 번호, 폰트주소
    text_box = slide.shapes.placeholders[cnt].text_frame
    text_box.fit_text(font_file=fontAdr)
    return text_box

#텍스트박스 객체에 새로운 텍스트를 추가한다. (줄단위)
def NewLine(textbox,text,font,size):#텍스트박스 객체, 내용, 폰트, 크기
    line = textbox.add_paragraph()
    line.font.name = font
    line.font.size = Pt(size)
    line.text = text

#ppt를 저장한다.
def write(prs,file_name):
    prs.save(file_name)




class TextData:
    def __init__(self, mainTitle, subTitle, midTitles, slideTitles, slideContents):
        self._mainTitle = mainTitle
        self._subTitle = subTitle
        self._midTitles = midTitles
        self._slideTitles = slideTitles
        self._slideContents = slideContents

    def __print__(self):
        print("---TextData Print---")
        print("제목: " + self._mainTitle)
        print("부제목: " + self._subTitle)
        print("중제목들 : ",end='')
        print(self._midTitles)
        print("슬라이드소제목들: ",end='')
        print(self._slideTitles)
        print("슬라이드내용들: ",end='')
        print(self._slideContents)

    def toJSON(self):
        return json.dumps(self, default=lambda o: o.__dict__, ensure_ascii = False, sort_keys=True, indent=4)

class PPTData:
    def __init__(self, textData):
        self._textData = textData


# 디폴트 타입
# [String] - 한줄내용
class SlideType:
    def __init__(self, lines):
        self._lines = lines
    def generate(self):
        print('generated')

# 타임라인 타입 (일정, 과정, 단계)
# [(String, String)] - (시간, 한줄내용)
class SlideType_timeline(SlideType):
    def __init__(self, timeTuples):
        self._timeTuples = timeTuples


# h5 타입 (비교대조 등)
# [(String,[String])] - (헤딩,[내용들])
class SlideType_h5(SlideType):
    def __init__(self, h5Tuples):
        self._h5Tuples = h5Tuples


# 정의 타입 (? -> "")
# String
class SlideType_definition(SlideType):
    def __init__(self, defStr):
        self._defStr = defStr
# def slide_test():
#     slides = [SlideType('default'),SlideType_h5([('SWM',['1','2','3']),('SWM',['1','2','3'])])]
# 
#     for Sample_sld in slides:
#         Sample_sld.generate()
